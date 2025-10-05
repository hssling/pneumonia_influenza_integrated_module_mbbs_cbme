import re
from pptx import Presentation
from pptx.util import Inches, Pt

def add_text_with_formatting(paragraph, text):
    # Split on ** for bold
    parts = re.split(r'(\*\*.+\*\*)', text)
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            run = paragraph.add_run()
            run.text = part[2:-2]
            run.font.bold = True
        else:
            run = paragraph.add_run()
            run.text = part
    # Similarly for *italic* if needed
    # For now, just bold

def parse_markdown_content(text):
    """Process markdown content to extract paragraphs, bullets, tables, etc."""
    lines = text.strip().split('\n')
    paragraphs = []
    bullets = []
    current_list = []
    table_rows = None

    for line in lines:
        line = line.strip()
        if not line:
            if table_rows:
                # End table
                paragraphs.append({'type': 'table', 'rows': table_rows})
                table_rows = None
            continue
        if line.count('|') > 1 and '|' in line:
            # Table row
            if table_rows is None:
                if bullets:
                    paragraphs.append({'type': 'bullets', 'items': bullets})
                    bullets = []
                table_rows = []
            table_rows.append([cell.strip() for cell in line.split('|') if cell.strip()])
        else:
            # End table if any
            if table_rows:
                paragraphs.append({'type': 'table', 'rows': table_rows})
                table_rows = None
            if line.startswith('### '):
                # Subsection
                if bullets:
                    paragraphs.append({'type': 'bullets', 'items': bullets})
                    bullets = []
                paragraphs.append({'type': 'heading', 'text': line[4:]})
            elif line.startswith('**') and line.endswith('**') and ':' in line:
                # Bold line like **Stats:**
                if bullets:
                    paragraphs.append({'type': 'bullets', 'items': bullets})
                    bullets = []
                paragraphs.append({'type': 'bold_text', 'text': line.strip('*')})
            elif line.startswith('- '):
                bullets.append(line[2:])
            elif line.startswith('*') and line.endswith('*'):
                if bullets:
                    current_list.append(line.strip('*'))
            else:
                # Regular text
                if bullets:
                    paragraphs.append({'type': 'bullets', 'items': bullets})
                    bullets = []
                paragraphs.append({'type': 'text', 'text': line})

    if bullets:
        paragraphs.append({'type': 'bullets', 'items': bullets})
    if table_rows:
        paragraphs.append({'type': 'table', 'rows': table_rows})

    return paragraphs

def create_pptx_from_md(md_file, pptx_file):
    # Read markdown
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split into slides
    slide_sections = content.split('---')
    prs = Presentation()

    for section in slide_sections:
        section = section.strip()
        if not section:
            continue

        lines = section.split('\n')
        slide_title = None
        slide_content = []

        for line in lines:
            if line.startswith('## Slide '):
                # Extract title after colon
                m = re.match(r'## Slide \d+: (.+)', line)
                if m:
                    slide_title = m.group(1)
                break

        if slide_title:
            slide_content = '\n'.join(lines[1:])  # Everything after the title line
        else:
            continue  # Skip sections without slide title

        # Create slide
        if slide_title and 'Title Slide' in slide_title:
            slide_layout = prs.slide_layouts[0]  # Title slide
        else:
            slide_layout = prs.slide_layouts[1]  # Title and content

        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = slide_title or 'Slide'

        if slide_content and len(slide.shapes.placeholders) > 1:
            content_placeholder = slide.shapes.placeholders[1]
            tf = content_placeholder.text_frame
            tf.clear()  # Clear any existing content

            paragraphs = parse_markdown_content(slide_content)
            for p in paragraphs:
                if p['type'] == 'bullets':
                    for item in p['items']:
                        p_tf = tf.add_paragraph()
                        add_text_with_formatting(p_tf, item.strip('*'))
                        p_tf.level = 0
                elif p['type'] == 'text':
                    p_tf = tf.add_paragraph()
                    add_text_with_formatting(p_tf, p['text'])
                elif p['type'] == 'heading':
                    p_tf = tf.add_paragraph()
                    p_tf.add_run()
                    run = p_tf.runs[0]
                    run.text = p['text']
                    run.font.bold = True
                    run.font.size = Pt(18)
                elif p['type'] == 'bold_text':
                    p_tf = tf.add_paragraph()
                    add_text_with_formatting(p_tf, p['text'])
                    for run in p_tf.runs:
                        run.font.bold = True
                elif p['type'] == 'table':
                    rows = p['rows']
                    num_rows = len(rows)
                    num_cols = max(len(row) for row in rows) if rows else 0
                    if num_cols > 0:
                        table = slide.shapes.add_table(num_rows, num_cols, Inches(1), Inches(3), Inches(9), Inches(3)).table
                        for i, row in enumerate(rows):
                            for j, cell in enumerate(row):
                                if j < num_cols:
                                    cell_object = table.cell(i, j)
                                    cell_object.text_frame.clear()
                                    if cell:
                                        p_tf = cell_object.text_frame.add_paragraph()
                                        add_text_with_formatting(p_tf, cell)

    # Save the presentation
    prs.save(pptx_file)
    print(f'PPTX created: {pptx_file}')

if __name__ == '__main__':
    import sys
    if len(sys.argv) != 3:
        print('Usage: python convert_md_to_pptx.py <md_file> <pptx_file>')
        sys.exit(1)

    md_file = sys.argv[1]
    pptx_file = sys.argv[2]
    create_pptx_from_md(md_file, pptx_file)
