#!/usr/bin/env python3
"""
Diabetes CBME Training Presentation Generator
Creates PowerPoint presentation (.pptx) from diabetes training materials
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

class DiabetesPresentationCreator:
    """
    Creates PowerPoint presentation from diabetes CBME training content
    """

    def __init__(self):
        self.presentation = Presentation()
        self.title_slide_layout = self.presentation.slide_layouts[0]
        self.content_slide_layout = self.presentation.slide_layouts[1]
        self.two_content_layout = self.presentation.slide_layouts[3]

        # Define color scheme
        self.colors = {
            'primary': RGBColor(0, 102, 204),    # Blue
            'secondary': RGBColor(0, 153, 51),   # Green
            'accent': RGBColor(255, 153, 0),     # Orange
            'text': RGBColor(50, 50, 50),        # Dark gray
            'white': RGBColor(255, 255, 255),    # White
            'light_gray': RGBColor(240, 240, 240) # Light gray
        }

    def add_title_slide(self):
        """Add title slide"""
        slide = self.presentation.slides.add_slide(self.title_slide_layout)

        # Title
        title = slide.shapes.title
        title.text = "Diabetes Mellitus Training"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        title.text_frame.paragraphs[0].font.bold = True

        # Subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = """CBME Module for MBBS Students
Competency-Based Medical Education

Dr. Siddalingaiah H S | MBBS, MD
Shridevi Institute of Medical Sciences & Research Hospital"""

        # Format subtitle
        for paragraph in subtitle.text_frame.paragraphs:
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = self.colors['text']

        return slide

    def add_overview_slide(self):
        """Add overview slide"""
        slide = self.presentation.slides.add_slide(self.content_slide_layout)

        # Title
        shapes = slide.shapes
        shapes.title.text = "Module Overview"
        shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
        shapes.title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']

        # Content
        content = """• Epidemiology & Global Prevalence of Diabetes
• Pathophysiology: Type 1 vs Type 2 Diabetes
• Clinical Features & Diagnosis (ADA 2023 Guidelines)
• Management & Treatment Strategies
• Acute & Chronic Complications
• Lifestyle Modification & Prevention
• Community Medicine & Screening Programs"""

        tf = shapes.placeholders[1].text_frame
        tf.text = content

        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.line_spacing = 1.2

        return slide

    def add_epidemiology_slide(self):
        """Add epidemiology slide"""
        slide = self.presentation.slides.add_slide(self.two_content_layout)

        # Title
        shapes = slide.shapes
        shapes.title.text = "Epidemiology & Global Impact"
        shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
        shapes.title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']

        # Left content - Statistics
        tf_left = shapes.placeholders[1].text_frame
        tf_left.text = "Global Diabetes Prevalence (IDF Atlas 2023)"

        p = tf_left.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True

        tf_left.add_paragraph().text = "• 537 million adults with diabetes worldwide"
        tf_left.add_paragraph().text = "• 1 in 10 adults affected"
        tf_left.add_paragraph().text = "• 75% live in low/middle-income countries"
        tf_left.add_paragraph().text = "• 6.7 million deaths annually"
        tf_left.add_paragraph().text = "• 966 billion USD economic impact"

        for paragraph in tf_left.paragraphs[1:]:
            paragraph.font.size = Pt(16)

        # Right content - India
        tf_right = shapes.placeholders[2].text_frame
        tf_right.text = "India's Diabetes Burden"

        p = tf_right.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['secondary']

        tf_right.add_paragraph().text = "• 77 million adults with diabetes"
        tf_right.add_paragraph().text = "• 25 million undiagnosed cases"
        tf_right.add_paragraph().text = "• 4th highest diabetes burden globally"
        tf_right.add_paragraph().text = "• Rural-urban divide: 4-6% rural vs 12-15% urban"
        tf_right.add_paragraph().text = "• Rising prevalence in younger population"

        for paragraph in tf_right.paragraphs[1:]:
            paragraph.font.size = Pt(16)

        return slide

    def add_pathophysiology_slide(self):
        """Add pathophysiology slide"""
        slide = self.presentation.slides.add_slide(self.two_content_layout)

        # Title
        shapes = slide.shapes
        shapes.title.text = "Pathophysiology: Type 1 vs Type 2 Diabetes"
        shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
        shapes.title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']

        # Left content - Type 1
        tf_left = shapes.placeholders[1].text_frame
        tf_left.text = "TYPE 1 DIABETES MELLITUS"

        p = tf_left.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 0, 0)  # Red

        tf_left.add_paragraph().text = "• Autoimmune destruction of β-cells"
        tf_left.add_paragraph().text = "• Complete insulin deficiency"
        tf_left.add_paragraph().text = "• Associated with HLA markers"
        tf_left.add_paragraph().text = "• Accounts for 5-10% of cases"
        tf_left.add_paragraph().text = "• Typically presents in childhood/adolescence"
        tf_left.add_paragraph().text = "• Requires lifetime insulin therapy"

        for paragraph in tf_left.paragraphs[1:]:
            paragraph.font.size = Pt(14)

        # Right content - Type 2
        tf_right = shapes.placeholders[2].text_frame
        tf_right.text = "TYPE 2 DIABETES MELLITUS"

        p = tf_right.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 123, 255)  # Blue

        tf_right.add_paragraph().text = "• Insulin resistance + relative insulin deficiency"
        tf_right.add_paragraph().text = "• Progressive β-cell dysfunction"
        tf_right.add_paragraph().text = "• Multifactorial genetics + environmental factors"
        tf_right.add_paragraph().text = "• Accounts for 90-95% of cases"
        tf_right.add_paragraph().text = "• Strong link with obesity and lifestyle"
        tf_right.add_paragraph().text = "• Can be managed with lifestyle + oral agents"

        for paragraph in tf_right.paragraphs[1:]:
            paragraph.font.size = Pt(14)

        return slide

    def add_diagnosis_slide(self):
        """Add diagnosis slide"""
        slide = self.presentation.slides.add_slide(self.content_slide_layout)

        # Title
        shapes = slide.shapes
        shapes.title.text = "Diagnosis & ADA 2023 Guidelines"
        shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
        shapes.title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']

        # Content
        content = """ADA DIAGNOSTIC CRITERIA:

Fasting Plasma Glucose (FPG):
• Normal: < 100 mg/dL (5.6 mmol/L)
• Prediabetes: 100-125 mg/dL (5.6-6.9 mmol/L)
• Diabetes: ≥ 126 mg/dL (7.0 mmol/L)

Oral Glucose Tolerance Test (OGTT):
• Normal: < 140 mg/dL (7.8 mmol/L) at 2 hours
• Prediabetes: 140-199 mg/dL (7.8-11.0 mmol/L) at 2 hours
• Diabetes: ≥ 200 mg/dL (11.1 mmol/L) at 2 hours

HbA1c (Hemoglobin A1c):
• Normal: < 5.7% (39 mmol/mol)
• Prediabetes: 5.7-6.4% (39-47 mmol/mol)
• Diabetes: ≥ 6.5% (48 mmol/mol)

Random Plasma Glucose:
• Diabetes: ≥ 200 mg/dL + symptoms"""

        tf = shapes.placeholders[1].text_frame
        tf.text = content

        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.line_spacing = 1.2

        # Bold key items
        for i, paragraph in enumerate(tf.paragraphs):
            if "• " in paragraph.text:
                paragraph.font.bold = True

        return slide

    def add_management_slide(self):
        """Add management slide"""
        slide = self.presentation.slides.add_slide(self.two_content_layout)

        # Title
        shapes = slide.shapes
        shapes.title.text = "Management & Treatment Strategies"
        shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
        shapes.title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']

        # Left content - Lifestyle + Oral Agents
        tf_left = shapes.placeholders[1].text_frame
        tf_left.text = "PHASE 1: LIFESTYLE + ORAL AGENTS"

        p = tf_left.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.colors['secondary']

        tf_left.add_paragraph().text = "Lifestyle Interventions:"
        tf_left.add_paragraph().text = "• Weight loss (5-10% body weight)"
        tf_left.add_paragraph().text = "• Mediterranean-style diet"
        tf_left.add_paragraph().text = "• Regular physical activity (150 min/week)"
        tf_left.add_paragraph().text = "• Smoking cessation"

        tf_left.add_paragraph().text = "Oral Hypoglycemic Agents:"
        tf_left.add_paragraph().text = "• Metformin (first-line)"
        tf_left.add_paragraph().text = "• SGLT2 inhibitors"
        tf_left.add_paragraph().text = "• DPP-4 inhibitors"
        tf_left.add_paragraph().text = "• GLP-1 receptor agonists"

        for paragraph in tf_left.paragraphs[1:]:
            paragraph.font.size = Pt(14)

        # Right content - Insulin Therapy
        tf_right = shapes.placeholders[2].text_frame
        tf_right.text = "PHASE 2: INSULIN THERAPY"

        p = tf_right.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 69, 0)  # Red-orange

        tf_right.add_paragraph().text = "Indications for Insulin:"
        tf_right.add_paragraph().text = "• HbA1c > 10% at diagnosis"
        tf_right.add_paragraph().text = "• Symptomatic hyperglycemia"
        tf_right.add_paragraph().text = "• Failure of oral agents + lifestyle"
        tf_right.add_paragraph().text = "• Pregnancy/GDM"
        tf_right.add_paragraph().text = "• Hospitalizations"

        tf_right.add_paragraph().text = "Insulin Regimens:"
        tf_right.add_paragraph().text = "• Basal insulin (Glargine, Detemir)"
        tf_right.add_paragraph().text = "• Premixed (70/30, 50/50)"
        tf_right.add_paragraph().text = "• Basal-bolus (MDI)"
        tf_right.add_paragraph().text = "• Insulin pump therapy"

        for paragraph in tf_right.paragraphs[1:]:
            paragraph.font.size = Pt(14)

        return slide

    def add_complications_slide(self):
        """Add complications slide"""
        slide = self.presentation.slides.add_slide(self.two_content_layout)

        # Title
        shapes = slide.shapes
        shapes.title.text = "Acute & Chronic Complications"
        shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
        shapes.title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']

        # Left content - Acute Complications
        tf_left = shapes.placeholders[1].text_frame
        tf_left.text = "ACUTE COMPLICATIONS"

        p = tf_left.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(220, 20, 60)  # Crimson

        tf_left.add_paragraph().text = "1. Hypoglycemia:"
        tf_left.add_paragraph().text = "   • Blood glucose < 70 mg/dL"
        tf_left.add_paragraph().text = "   • Treatment: 15-20g fast-acting carbs"
        tf_left.add_paragraph().text = "   • Severe cases: Glucagon/IM glucose"

        tf_left.add_paragraph().text = "2. Diabetic Ketoacidosis (DKA):"
        tf_left.add_paragraph().text = "   • Hyperglycemia + ketones + acidosis"
        tf_left.add_paragraph().text = "   • Treatment: Fluid resuscitation + insulin"

        tf_left.add_paragraph().text = "3. Hyperosmolar Hyperglycemic State (HHS):"
        tf_left.add_paragraph().text = "   • Severe hyperglycemia without ketoacidosis"
        tf_left.add_paragraph().text = "   • Treatment: Aggressive fluid therapy"

        for paragraph in tf_left.paragraphs[1:]:
            paragraph.font.size = Pt(12)

        # Right content - Chronic Complications
        tf_right = shapes.placeholders[2].text_frame
        tf_right.text = "CHRONIC MICROVASCULAR COMPLICATIONS"

        p = tf_right.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(139, 69, 19)  # Saddle brown

        tf_right.add_paragraph().text = "1. Diabetic Retinopathy:"
        tf_right.add_paragraph().text = "   • Leading cause of blindness"
        tf_right.add_paragraph().text = "   • Annual screening recommended"
        tf_right.add_paragraph().text = "   • Laser photocoagulation effective"

        tf_right.add_paragraph().text = "2. Diabetic Nephropathy:"
        tf_right.add_paragraph().text = "   • Microalbuminuria → proteinuria → ESRD"
        tf_right.add_paragraph().text = "   • Prevention: RAAS blockade"
        tf_right.add_paragraph().text = "   • 40% of ESRD cases in India"

        tf_right.add_paragraph().text = "3. Diabetic Neuropathy:"
        tf_right.add_paragraph().text = "   • Sensory + autonomic involvement"
        tf_right.add_paragraph().text = "   • Prevention: Glycemic control"
        tf_right.add_paragraph().text = "   • Most common chronic complication"

        for paragraph in tf_right.paragraphs[1:]:
            paragraph.font.size = Pt(12)

        return slide

    def add_prevention_slide(self):
        """Add prevention slide"""
        slide = self.presentation.slides.add_slide(self.content_slide_layout)

        # Title
        shapes = slide.shapes
        shapes.title.text = "Prevention & Lifestyle Management"
        shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
        shapes.title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']

        # Content
        content = """PRIMARY PREVENTION STRATEGIES:

Lifestyle Interventions:
• Mediterranean diet with high fiber intake
• Regular physical activity (brisk walking 30 min/day)
• Maintenance of healthy BMI (18.5-24.9 kg/m²)
• Smoking cessation and alcohol moderation
• Stress management and adequate sleep

Population-Level Interventions:
• Sugar-sweetened beverage taxation
• Labeling requirements for packaged foods
• Workplace wellness programs
• Community screening campaigns
• Public awareness campaigns

SECONDARY PREVENTION:
• Regular screening for high-risk groups
• Early diagnosis and treatment
• Cardiovascular risk factor management
• Annual eye and foot examinations
• HbA1c monitoring every 3-6 months

Evidence-Based Benefits:
• 58% reduction in diabetes incidence (DPP Trial)
• 20-30% CVD risk reduction
• Improved quality of life and life expectancy"""

        tf = shapes.placeholders[1].text_frame
        tf.text = content

        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.line_spacing = 1.2

        # Color key items
        for i, paragraph in enumerate(tf.paragraphs):
            if "•" in paragraph.text or ":" in paragraph.text:
                paragraph.font.bold = True

        return slide

    def add_case_study_slide(self):
        """Add case study slide"""
        slide = self.presentation.slides.add_slide(self.two_content_layout)

        # Title
        shapes = slide.shapes
        shapes.title.text = "CBME Case Study: Type 2 Diabetes Management"
        shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
        shapes.title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']

        # Left content - Case Presentation
        tf_left = shapes.placeholders[1].text_frame
        tf_left.text = "CASE PRESENTATION"

        p = tf_left.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True

        tf_left.add_paragraph().text = "45-year-old male, overweight (BMI 32)"
        tf_left.add_paragraph().text = "Chief Complaints:"
        tf_left.add_paragraph().text = "• Increased thirst and urination"
        tf_left.add_paragraph().text = "• Unexplained weight loss (8 kg)"
        tf_left.add_paragraph().text = "• Blurred vision"
        tf_left.add_paragraph().text = "• Fatigue"

        tf_left.add_paragraph().text = "Past History:"
        tf_left.add_paragraph().text = "• Hypertension (on medication)"
        tf_left.add_paragraph().text = "• Sedentary lifestyle"
        tf_left.add_paragraph().text = "• Family history: Mother with T2DM"

        tf_left.add_paragraph().text = "Laboratory Results:"
        tf_left.add_paragraph().text = "• RBS: 280 mg/dL"
        tf_left.add_paragraph().text = "• HbA1c: 9.2%"
        tf_left.add_paragraph().text = "• Serum creatinine: 1.1 mg/dL"

        for paragraph in tf_left.paragraphs[1:]:
            paragraph.font.size = Pt(14)

        # Right content - Management Plan
        tf_right = shapes.placeholders[2].text_frame
        tf_right.text = "CBME MANAGEMENT PLAN"

        p = tf_right.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.colors['secondary']

        tf_right.add_paragraph().text = "1. Patient Care:"
        tf_right.add_paragraph().text = "   • Lifestyle counseling"
        tf_right.add_paragraph().text = "   • SMBG education"
        tf_right.add_paragraph().text = "   • Hypoglycemia awareness"

        tf_right.add_paragraph().text = "2. Medical Knowledge:"
        tf_right.add_paragraph().text = "   • Metformin 500mg BD"
        tf_right.add_paragraph().text = "   • ACE inhibitor for hypertension"
        tf_right.add_paragraph().text = "   • Statin therapy"

        tf_right.add_paragraph().text = "3. Practice-Based Learning:"
        tf_right.add_paragraph().text = "   • Monthly follow-ups"
        tf_right.add_paragraph().text = "   • HbA1c target: <7.0%"
        tf_right.add_paragraph().text = "   • Patient education materials"

        tf_right.add_paragraph().text = "4. Communication Skills:"
        tf_right.add_paragraph().text = "   • Family counseling"
        tf_right.add_paragraph().text = "   • Support group referral"

        for paragraph in tf_right.paragraphs[1:]:
            paragraph.font.size = Pt(14)

        return slide

    def add_assessment_slide(self):
        """Add assessment slide"""
        slide = self.presentation.slides.add_slide(self.content_slide_layout)

        # Title
        shapes = slide.shapes
        shapes.title.text = "CBME Assessment & Learning Objectives"
        shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
        shapes.title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']

        # Content
        content = """CBME COMPETENCIES (MBBS LEVEL):

1. Patient Care:
   • Diagnose diabetes mellitus using clinical criteria
   • Initiate appropriate treatment plans
   • Monitor glycemic control and complications

2. Medical Knowledge:
   • Understand pathophysiology of hyperglycemia
   • Know pharmacological and non-pharmacological treatments
   • Recognize complications and their management

3. Practice-Based Learning:
   • Use evidence-based guidelines (ADA 2023)
   • Interpret laboratory results appropriately
   • Make rational clinical decisions

4. Communication Skills:
   • Educate patients about diabetes self-management
   • Counsel on lifestyle modifications
   • Coordinate with multidisciplinary team

5. Professionalism:
   • Maintain patient confidentiality
   • Provide culturally competent care
   • Uphold ethical standards

6. Systems-Based Practice:
   • Utilize community resources for diabetes care
   • Understand healthcare delivery systems
   • Advocate for preventive measures"""

        tf = shapes.placeholders[1].text_frame
        tf.text = content

        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.line_spacing = 1.2

        # Bold competency titles
        for paragraph in tf.paragraphs:
            if "1. " in paragraph.text or "2. " in paragraph.text or "3. " in paragraph.text or "4. " in paragraph.text or "5. " in paragraph.text or "6. " in paragraph.text:
                paragraph.font.bold = True
                paragraph.font.size = Pt(16)

        return slide

    def add_resources_slide(self):
        """Add resources slide"""
        slide = self.presentation.slides.add_slide(self.two_content_layout)

        # Title
        shapes = slide.shapes
        shapes.title.text = "Educational Resources & References"
        shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
        shapes.title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']

        # Left content - Books & Guidelines
        tf_left = shapes.placeholders[1].text_frame
        tf_left.text = "KEY REFERENCES"

        p = tf_left.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True

        tf_left.add_paragraph().text = "American Diabetes Association:"
        tf_left.add_paragraph().text = "• Standards of Care 2023"
        tf_left.add_paragraph().text = "• Complete Type 2 Diabetes Guideline"

        tf_left.add_paragraph().text = "International Diabetes Federation:"
        tf_left.add_paragraph().text = "• IDF Diabetes Atlas 2023"
        tf_left.add_paragraph().text = "• Regional reports"

        tf_left.add_paragraph().text = "Williams Textbook of Endocrinology"
        tf_left.add_paragraph().text = "Harrison's Principles of Internal Medicine"
        tf_left.add_paragraph().text = "Joslin's Diabetes Mellitus"

        for paragraph in tf_left.paragraphs[1:]:
            paragraph.font.size = Pt(14)

        # Right content - Online Resources
        tf_right = shapes.placeholders[2].text_frame
        tf_right.text = "ONLINE RESOURCES"

        p = tf_right.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.colors['accent']

        tf_right.add_paragraph().text = "Professional Websites:"
        tf_right.add_paragraph().text = "• diabetes.org (ADA)"
        tf_right.add_paragraph().text = "• idf.org (IDF)"
        tf_right.add_paragraph().text = "• who.int/diabetes"

        tf_right.add_paragraph().text = "Educational Platforms:"
        tf_right.add_paragraph().text = "• Coursera - Diabetes courses"
        tf_right.add_paragraph().text = "• edX - Endocrinology programs"
        tf_right.add_paragraph().text = "• Khan Academy - Basic physiology"

        tf_right.add_paragraph().text = "Latest Research:"
        tf_right.add_paragraph().text = "• NEJM Diabetes reviews"
        tf_right.add_paragraph().text = "• The Lancet Diabetes & Endocrinology"
        tf_right.add_paragraph().text = "• PubMed diabetes literature"

        for paragraph in tf_right.paragraphs[1:]:
            paragraph.font.size = Pt(14)

        return slide

    def add_summary_slide(self):
        """Add summary slide"""
        slide = self.presentation.slides.add_slide(self.content_slide_layout)

        # Title
        shapes = slide.shapes
        shapes.title.text = "Summary: Key Takeaways"
        shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
        shapes.title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']

        # Content
        content = """DIABETES MELLITUS - CRITICAL CONCEPTS:

🔹 Global epidemic affecting 537 million people
🔹 Type 2 diabetes accounts for 90-95% of cases
🔹 Diagnosis requires HbA1c ≥ 6.5% or confirmatory lab results
🔹 Metformin remains first-line therapy for T2DM
🔹 Glycemic control prevents microvascular complications
🔹 Lifestyle modification can prevent diabetes development
🔹 Regular screening advised for high-risk populations

CLINICAL DECISION POINTS:
• RBS > 200 mg/dL + symptoms = diabetes
• HbA1c > 10% at diagnosis = insulin initiation
• Chronic hyperglycemia = aggressive risk factor management
• Patient education = cornerstone of diabetes care

CBME FOCUS:
• Physician competence in diagnosis and management
• Evidence-based therapeutic decision making
• Patient-centered communication and education
• Lifelong learning and quality improvement

REMEMBER: Diabetes is a preventable disease. Early intervention saves lives and reduces healthcare burden."""

        tf = shapes.placeholders[1].text_frame
        tf.text = content

        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.line_spacing = 1.2

        # Color emojis
        for paragraph in tf.paragraphs:
            if "🔹" in paragraph.text or "•" in paragraph.text:
                paragraph.font.bold = True

        return slide

    def save_presentation(self, filename="diabetes_cbme_training.pptx"):
        """Save the PowerPoint presentation"""
        try:
            self.presentation.save(filename)
            print(f"✅ PowerPoint presentation saved as '{filename}'")
            return filename
        except Exception as e:
            print(f"❌ Error saving presentation: {e}")
            return None

    def create_full_presentation(self):
        """Create the complete PowerPoint presentation"""
        print("🩺 Creating Diabetes CBME Training Presentation...")

        # Add all slides in logical order
        self.add_title_slide()
        self.add_overview_slide()
        self.add_epidemiology_slide()
        self.add_pathophysiology_slide()
        self.add_diagnosis_slide()
        self.add_management_slide()
        self.add_complications_slide()
        self.add_prevention_slide()
        self.add_case_study_slide()
        self.add_assessment_slide()
        self.add_resources_slide()
        self.add_summary_slide()

        print(f"✅ Created presentation with {len(self.presentation.slides)} slides")
        return self.save_presentation()


def main():
    """Main function to generate diabetes training PowerPoint"""
    print("=" * 60)
    print("🩺 DIABETES CBME TRAINING PRESENTATION GENERATOR")
    print("=" * 60)
    print()
    print("This script creates a comprehensive PowerPoint presentation")
    print("covering diabetes mellitus CBME training for MBBS students.")
    print()
    print("Features:")
    print("• Complete module overview and epidemiology")
    print("• Pathophysiology and clinical diagnosis")
    print("• Management strategies and treatment guidelines")
    print("• Complications and prevention approaches")
    print("• CBME case study and assessment criteria")
    print("• Educational resources and key references")
    print()
    print("Requirements:")
    print("pip install python-pptx")
    print()

    try:
        creator = DiabetesPresentationCreator()
        filename = creator.create_full_presentation()

        if filename:
            print("\n🎉 SUCCESS!")
            print(f"📊 Presentation saved as: {filename}")
            print(f"📈 Total slides: 12")
            print("📚 Ready for MBBS diabetes training sessions!")

            print("\n📎 Next Steps:")
            print("• Open the .pptx file in PowerPoint")
            print("• Review and customize content as needed")
            print("• Add images/diagrams from medical literature")

    except Exception as e:
        print(f"❌ Error: {e}")
        print("💡 Make sure you have python-pptx installed: pip install python-pptx")

if __name__ == "__main__":
    main()
